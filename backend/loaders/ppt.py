"""
PowerPoint Loader Module

Extracts text from PPTX presentations.
"""

from pptx import Presentation
from langchain_core.documents import Document
from typing import List


def load_ppt(file_path: str) -> List[Document]:
    """
    Load a PPTX file and return a list of Documents (one per slide).
    
    Each slide becomes a separate document for granular retrieval.
    """
    documents = []
    
    try:
        prs = Presentation(file_path)
        print(f"📊 Loading PPTX: {file_path} ({len(prs.slides)} slides)")
        
        for i, slide in enumerate(prs.slides):
            text_runs = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            text_runs.append(run.text)
            
            text = "\n".join(text_runs).strip()
            
            if text:
                documents.append(Document(
                    page_content=text,
                    metadata={
                        "source": file_path,
                        "slide": i + 1,
                        "total_slides": len(prs.slides)
                    }
                ))
                
        print(f"   ✅ Extracted {len(documents)} slides")
        
    except Exception as e:
        print(f"   ❌ Error loading PPTX {file_path}: {e}")
        
    return documents
